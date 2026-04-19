"""
Generate Tether — ULM GDSC Hawkathon presentation.
Run: python3 make_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ───────────────────────────────────────────────────────────────────
ORANGE  = RGBColor(0xE8, 0x93, 0x0A)
DARK    = RGBColor(0x1C, 0x14, 0x10)
CREAM   = RGBColor(0xFA, 0xFA, 0xF7)
SAND    = RGBColor(0xE2, 0xD4, 0xBC)
MUTED   = RGBColor(0x9C, 0x8A, 0x72)
GREEN   = RGBColor(0x15, 0xA8, 0x7A)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL= RGBColor(0x2C, 0x22, 0x1A)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


def add_rect(slide, x, y, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def add_text(slide, text, x, y, w, h,
             size=24, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=16, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, space_before=None):
    p   = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return p


def cream_bg(slide):
    add_rect(slide, 0, 0, W, H, CREAM)


def accent_bar(slide, height=Inches(0.08)):
    add_rect(slide, 0, 0, W, height, ORANGE)


def tag(slide, label, x, y, bg=ORANGE, fg=WHITE, size=11):
    w, h = Inches(1.8), Inches(0.32)
    add_rect(slide, x, y, w, h, bg)
    add_text(slide, label, x, y, w, h, size=size, bold=True,
             color=fg, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
s = prs.slides.add_slide(BLANK)
# Full dark background
add_rect(s, 0, 0, W, H, DARK)
# Orange accent top strip
add_rect(s, 0, 0, W, Inches(0.12), ORANGE)
# Orange side accent
add_rect(s, 0, 0, Inches(0.12), H, ORANGE)
# Decorative circle (large, subtle)
add_rect(s, Inches(8.5), Inches(1.0), Inches(5), Inches(5), CHARCOAL)

# Hawkathon badge
tag(s, "ULM GDSC HAWKATHON 2025", Inches(0.5), Inches(1.2), bg=ORANGE, fg=WHITE, size=10)

# Title
add_text(s, "Tether", Inches(0.5), Inches(1.8), Inches(7), Inches(1.6),
         size=96, bold=True, color=WHITE)
# Subtitle
add_text(s, "AI-Powered Career Agent for Monroe, Louisiana",
         Inches(0.5), Inches(3.3), Inches(7.5), Inches(0.8),
         size=22, bold=False, color=SAND)
# Tagline
add_text(s, "Find jobs. Reach recruiters. Land the role.",
         Inches(0.5), Inches(4.0), Inches(7.5), Inches(0.6),
         size=16, bold=False, color=MUTED)

# Divider
add_rect(s, Inches(0.5), Inches(4.9), Inches(3), Inches(0.04), ORANGE)

# Team
add_text(s, "Team RaiCutz",
         Inches(0.5), Inches(5.1), Inches(6), Inches(0.45),
         size=18, bold=True, color=ORANGE)
add_text(s, "George Khawas  ·  Arhyael Zoaka  ·  Aashutosh Khanal  ·  Aayush Rai",
         Inches(0.5), Inches(5.55), Inches(9), Inches(0.4),
         size=13, bold=False, color=SAND)

# =============================================================================
# SLIDE 2 — The Problem
# =============================================================================
s = prs.slides.add_slide(BLANK)
cream_bg(s)
accent_bar(s)

add_text(s, "The Problem", Inches(0.7), Inches(0.35), Inches(8), Inches(0.6),
         size=11, bold=True, color=ORANGE)
add_text(s, "ULM students face a broken job search",
         Inches(0.7), Inches(0.9), Inches(10), Inches(0.9),
         size=38, bold=True, color=DARK)

problems = [
    ("🔍", "Hard to find",    "Local Monroe-area job listings are scattered across dozens of sites — no single source of truth."),
    ("📭", "No connections",  "Students don't know which recruiter or HR contact to reach at a company, so cold outreach never happens."),
    ("✍️", "Blank page panic", "Even motivated students stall when asked to write a personalized cold email from scratch."),
    ("⏱️", "Time drain",      "Manually searching, drafting, and sending emails for 10+ companies takes days — most students give up."),
]

cols = [Inches(0.7), Inches(3.6), Inches(6.5), Inches(9.4)]
for i, (icon, title, body) in enumerate(problems):
    x = cols[i]
    add_rect(s, x, Inches(2.2), Inches(2.6), Inches(3.8), SAND)
    add_text(s, icon,  x + Inches(0.15), Inches(2.35), Inches(0.6), Inches(0.6), size=26)
    add_text(s, title, x + Inches(0.15), Inches(3.0),  Inches(2.3), Inches(0.5),
             size=14, bold=True, color=DARK)
    add_text(s, body,  x + Inches(0.15), Inches(3.5),  Inches(2.3), Inches(2.0),
             size=11, bold=False, color=MUTED)

# =============================================================================
# SLIDE 3 — Our Solution
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, DARK)
add_rect(s, 0, 0, W, Inches(0.12), ORANGE)

add_text(s, "Our Solution", Inches(0.7), Inches(0.35), Inches(8), Inches(0.5),
         size=11, bold=True, color=ORANGE)
add_text(s, "Meet Tether", Inches(0.7), Inches(0.85), Inches(9), Inches(0.9),
         size=48, bold=True, color=WHITE)
add_text(s, "An end-to-end AI career agent that automates the entire job-search pipeline — "
            "from resume analysis to emails landing in recruiter inboxes.",
         Inches(0.7), Inches(1.75), Inches(9), Inches(0.9),
         size=16, bold=False, color=SAND)

steps = [
    ("01", "Upload Resume",     "Paste or upload your PDF. Claude reads it and builds your career profile instantly."),
    ("02", "Discover Jobs",     "AI scrapes Indeed, LinkedIn & ZipRecruiter for Monroe-area roles matching your background."),
    ("03", "Find Contacts",     "Tether surfaces HR managers and recruiters at every company — no cold searching required."),
    ("04", "Draft Emails",      "Claude writes a personalized cold email for each contact using your real resume as context."),
    ("05", "Auto-Send",         "One click — all emails go out from your Gmail. Real names. Real signatures. Real results."),
]

for i, (num, title, body) in enumerate(steps):
    x = Inches(0.7) + i * Inches(2.5)
    add_rect(s, x, Inches(3.0), Inches(2.3), Inches(3.8), CHARCOAL)
    add_text(s, num,   x + Inches(0.15), Inches(3.1),  Inches(0.8), Inches(0.5),
             size=28, bold=True, color=ORANGE)
    add_text(s, title, x + Inches(0.15), Inches(3.65), Inches(2.0), Inches(0.5),
             size=13, bold=True, color=WHITE)
    add_text(s, body,  x + Inches(0.15), Inches(4.2),  Inches(2.0), Inches(2.4),
             size=11, bold=False, color=SAND)

# =============================================================================
# SLIDE 4 — Key Features
# =============================================================================
s = prs.slides.add_slide(BLANK)
cream_bg(s)
accent_bar(s)

add_text(s, "Features", Inches(0.7), Inches(0.35), Inches(8), Inches(0.5),
         size=11, bold=True, color=ORANGE)
add_text(s, "Everything a student needs to land the job",
         Inches(0.7), Inches(0.85), Inches(10), Inches(0.8),
         size=36, bold=True, color=DARK)

features = [
    ("📄 Resume Analysis",      "Claude reads your resume and extracts skills, experience, and gaps — building a rich profile used across all features."),
    ("🗺️ Career Report",         "Get a personalized 90-day action plan, skill gap analysis, and employer shortlist based on your background."),
    ("💼 Live Job Feed",         "Real-time job scraping from Indeed, LinkedIn, and ZipRecruiter — filtered for Monroe, LA and nearby cities."),
    ("🎯 Job-Specific Plans",    "Click any listing and generate a tailored 90-day plan that references the company and role by name."),
    ("📇 Contact Generation",   "HR managers and recruiters are surfaced for every company — with role-based generic fallback emails."),
    ("⚡ One-Click Automator",   "Enter job title + Gmail credentials → Tether searches jobs, drafts emails, and sends them automatically."),
    ("✉️ AI Email Drafts",       "Every cold email is written by Claude using your actual resume context, the job description, and your name."),
    ("🤖 Career AI Chat",        "Ask anything — salary negotiation, interview prep, resume feedback — answered with your profile in mind."),
]

for i, (title, body) in enumerate(features):
    row, col = divmod(i, 4)
    x = Inches(0.7) + col * Inches(3.15)
    y = Inches(2.0) + row * Inches(2.3)
    add_rect(s, x, y, Inches(3.0), Inches(2.1), SAND)
    add_text(s, title, x + Inches(0.15), y + Inches(0.12), Inches(2.7), Inches(0.5),
             size=12, bold=True, color=DARK)
    add_text(s, body,  x + Inches(0.15), y + Inches(0.6),  Inches(2.7), Inches(1.4),
             size=10, bold=False, color=MUTED)

# =============================================================================
# SLIDE 5 — Tech Stack
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, DARK)
add_rect(s, 0, 0, W, Inches(0.12), ORANGE)

add_text(s, "Tech Stack", Inches(0.7), Inches(0.35), Inches(8), Inches(0.5),
         size=11, bold=True, color=ORANGE)
add_text(s, "Built with modern, production-grade tools",
         Inches(0.7), Inches(0.85), Inches(10), Inches(0.8),
         size=36, bold=True, color=WHITE)

categories = [
    ("Frontend",  GREEN,  ["Next.js 14 (App Router)", "TypeScript", "Tailwind CSS", "Shadcn/UI", "Axios"]),
    ("Backend",   ORANGE, ["FastAPI (Python)", "Uvicorn", "aiosqlite / SQLite", "SlowAPI rate limiting", "JWT auth + Fernet encryption"]),
    ("AI / Data", RGBColor(0x60,0x8B,0xFF), ["Anthropic Claude Sonnet 4.6", "python-jobspy (scraping)", "pypdf (resume parsing)", "Deterministic contact gen"]),
    ("Infra",     SAND,   ["Gmail SMTP (App Passwords)", "Cookie-based sessions", "CORS hardened", "Env-var secrets management"]),
]

for i, (cat, color, items) in enumerate(categories):
    x = Inches(0.7) + i * Inches(3.15)
    add_rect(s, x, Inches(2.1), Inches(3.0), Inches(0.04), color)
    add_text(s, cat, x, Inches(2.2), Inches(3.0), Inches(0.5),
             size=14, bold=True, color=color)
    for j, item in enumerate(items):
        add_text(s, f"· {item}", x, Inches(2.8) + j * Inches(0.55), Inches(3.0), Inches(0.5),
                 size=12, bold=False, color=SAND)

# =============================================================================
# SLIDE 6 — Architecture
# =============================================================================
s = prs.slides.add_slide(BLANK)
cream_bg(s)
accent_bar(s)

add_text(s, "Architecture", Inches(0.7), Inches(0.35), Inches(8), Inches(0.5),
         size=11, bold=True, color=ORANGE)
add_text(s, "How Tether works under the hood",
         Inches(0.7), Inches(0.85), Inches(10), Inches(0.8),
         size=36, bold=True, color=DARK)

boxes = [
    (Inches(0.5),  Inches(2.3), "Next.js Frontend",   "Dashboard · Automator\nChat · Report Tabs",    SAND,    DARK),
    (Inches(3.8),  Inches(2.3), "FastAPI Backend",     "8 routers · JWT auth\nRate limiting · CORS",   DARK,    WHITE),
    (Inches(7.1),  Inches(2.3), "Claude Sonnet 4.6",   "Resume analysis\nEmail drafts · Chat · Plans", ORANGE,  WHITE),
    (Inches(10.4), Inches(2.3), "SQLite DB",           "Users · Reports\nDrafts · Contacts · Plans",   GREEN,   WHITE),
    (Inches(0.5),  Inches(4.9), "JobSpy Scraper",      "Indeed · LinkedIn\nZipRecruiter",              SAND,    DARK),
    (Inches(3.8),  Inches(4.9), "Contact Generator",   "Deterministic seeded\nHR contact synthesis",   SAND,    DARK),
    (Inches(7.1),  Inches(4.9), "Gmail SMTP",          "Encrypted App Passwords\nAuto outreach send",  SAND,    DARK),
    (Inches(10.4), Inches(4.9), "pypdf Parser",        "Resume text extraction\nContext for Claude",    SAND,    DARK),
]

for x, y, title, body, bg, fg in boxes:
    add_rect(s, x, y, Inches(2.6), Inches(1.7), bg)
    add_text(s, title, x + Inches(0.15), y + Inches(0.12), Inches(2.3), Inches(0.45),
             size=12, bold=True, color=fg)
    add_text(s, body,  x + Inches(0.15), y + Inches(0.6),  Inches(2.3), Inches(1.0),
             size=10, bold=False, color=fg if fg == WHITE else MUTED)

# Arrow labels
for x in [Inches(3.2), Inches(6.5), Inches(9.8)]:
    add_text(s, "→", x, Inches(2.9), Inches(0.5), Inches(0.5),
             size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 7 — Impact & Vision
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, DARK)
add_rect(s, 0, 0, W, Inches(0.12), ORANGE)

add_text(s, "Impact & Vision", Inches(0.7), Inches(0.35), Inches(8), Inches(0.5),
         size=11, bold=True, color=ORANGE)
add_text(s, "Why Tether matters",
         Inches(0.7), Inches(0.85), Inches(9), Inches(0.8),
         size=42, bold=True, color=WHITE)

stats = [
    ("10×", "faster outreach", "What took days of manual work now takes one click."),
    ("100%", "local focus",    "Tuned specifically for Monroe, LA — not generic job boards."),
    ("0$",  "to get started",  "Free to run locally. No subscription. No paywalled contacts."),
]

for i, (num, label, desc) in enumerate(stats):
    x = Inches(0.7) + i * Inches(4.2)
    add_rect(s, x, Inches(2.2), Inches(3.8), Inches(2.2), CHARCOAL)
    add_text(s, num,   x + Inches(0.2), Inches(2.35), Inches(3.4), Inches(0.85),
             size=54, bold=True, color=ORANGE)
    add_text(s, label, x + Inches(0.2), Inches(3.2),  Inches(3.4), Inches(0.4),
             size=16, bold=True, color=WHITE)
    add_text(s, desc,  x + Inches(0.2), Inches(3.65), Inches(3.4), Inches(0.6),
             size=11, bold=False, color=SAND)

add_text(s, "What's next",
         Inches(0.7), Inches(5.0), Inches(5), Inches(0.45),
         size=14, bold=True, color=ORANGE)

next_items = [
    "LinkedIn OAuth for real recruiter data",
    "SMS / WhatsApp follow-up automation",
    "Interview scheduler with calendar integration",
    "Expand beyond Monroe to all of Louisiana",
]
for i, item in enumerate(next_items):
    add_text(s, f"→  {item}",
             Inches(0.7) + (i % 2) * Inches(6.3),
             Inches(5.5) + (i // 2) * Inches(0.5),
             Inches(6.0), Inches(0.45),
             size=12, bold=False, color=SAND)

# =============================================================================
# SLIDE 8 — Team
# =============================================================================
s = prs.slides.add_slide(BLANK)
cream_bg(s)
add_rect(s, 0, 0, W, Inches(0.12), ORANGE)

add_text(s, "The Team", Inches(0.7), Inches(0.35), Inches(8), Inches(0.5),
         size=11, bold=True, color=ORANGE)
add_text(s, "Team RaiCutz",
         Inches(0.7), Inches(0.85), Inches(9), Inches(0.9),
         size=48, bold=True, color=DARK)
add_text(s, "University of Louisiana Monroe  ·  ULM GDSC Hawkathon 2025",
         Inches(0.7), Inches(1.75), Inches(10), Inches(0.5),
         size=14, bold=False, color=MUTED)

members = [
    ("George Khawas",     "Full-Stack Development",  "Backend architecture,\nFastAPI routers, database design"),
    ("Arhyael Zoaka",     "AI & Automation",         "Claude integration,\nautomator pipeline, email systems"),
    ("Aashutosh Khanal",  "Frontend Development",    "Next.js UI, dashboard,\nReportTabs & chat interface"),
    ("Aayush Rai",        "Data & Job Search",       "JobSpy integration,\ncontact generation, scraping"),
]

for i, (name, role, desc) in enumerate(members):
    x = Inches(0.7) + i * Inches(3.15)
    add_rect(s, x, Inches(2.6), Inches(2.9), Inches(3.8), DARK)
    # Avatar circle placeholder
    add_rect(s, x + Inches(1.05), Inches(2.75), Inches(0.8), Inches(0.8), ORANGE)
    add_text(s, name[0], x + Inches(1.05), Inches(2.75), Inches(0.8), Inches(0.8),
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, name, x + Inches(0.15), Inches(3.7), Inches(2.6), Inches(0.5),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x + Inches(0.5), Inches(4.25), Inches(1.9), Inches(0.28), ORANGE)
    add_text(s, role, x + Inches(0.5), Inches(4.25), Inches(1.9), Inches(0.28),
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, x + Inches(0.15), Inches(4.65), Inches(2.6), Inches(1.5),
             size=10, bold=False, color=SAND, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 9 — Thank You / Closing
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, DARK)
add_rect(s, 0, 0, W, Inches(0.12), ORANGE)
add_rect(s, 0, H - Inches(0.12), W, Inches(0.12), ORANGE)

add_text(s, "Thank You", Inches(0), Inches(1.8), W, Inches(1.6),
         size=96, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Tether — Connecting ULM students to Monroe opportunities",
         Inches(0), Inches(3.5), W, Inches(0.7),
         size=18, bold=False, color=SAND, align=PP_ALIGN.CENTER)

add_rect(s, Inches(5.4), Inches(4.4), Inches(2.5), Inches(0.04), ORANGE)

add_text(s, "Team RaiCutz  ·  ULM GDSC Hawkathon 2025",
         Inches(0), Inches(4.6), W, Inches(0.5),
         size=13, bold=False, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, "George Khawas  ·  Arhyael Zoaka  ·  Aashutosh Khanal  ·  Aayush Rai",
         Inches(0), Inches(5.1), W, Inches(0.4),
         size=12, bold=False, color=SAND, align=PP_ALIGN.CENTER)

# =============================================================================
out = "/Users/rhe_pc/Tether/Tether_Hawkathon_RaiCutz.pptx"
prs.save(out)
print(f"Saved → {out}")
